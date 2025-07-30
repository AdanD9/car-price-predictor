#!/usr/bin/env python3
"""
Car Price Predictor - PowerPoint Presentation Generator
Creates a professional PowerPoint presentation from the technical documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

class PowerPointGenerator:
    def __init__(self, filename="Car_Price_Predictor_Presentation.pptx"):
        self.filename = filename
        self.prs = Presentation()
        
        # Define color scheme
        self.primary_color = RGBColor(37, 99, 235)  # Blue
        self.secondary_color = RGBColor(55, 65, 81)  # Dark gray
        self.accent_color = RGBColor(16, 185, 129)   # Green
        self.text_color = RGBColor(31, 41, 55)       # Dark text

    def add_title_slide(self):
        """Add title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Car Price Predictor"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        subtitle.text = "AI-Powered Automotive Valuation Platform\n\nTechnical Overview & Business Impact\nJuly 29, 2025"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.secondary_color

    def add_executive_summary(self):
        """Add executive summary slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🎯 Mission"
        
        p = tf.add_paragraph()
        p.text = "Revolutionize car valuation with AI-powered precision and instant results"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "📊 Key Metrics"
        
        p = tf.add_paragraph()
        p.text = "• 98.3% Prediction Accuracy"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Instant Results (< 3 seconds)"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 87+ ML Features Analyzed"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Production Ready on CarUbuntu Server"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "💼 Business Value"
        
        p = tf.add_paragraph()
        p.text = "Scalable SaaS revenue model with enterprise-grade security"
        p.level = 1

    def add_platform_overview(self):
        """Add platform overview slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Platform Overview"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🚗 What We Do"
        
        p = tf.add_paragraph()
        p.text = "Transform complex vehicle data into accurate price predictions using advanced machine learning"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "🎯 Target Users"
        
        p = tf.add_paragraph()
        p.text = "• Individual Buyers/Sellers: Personal vehicle valuation"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Automotive Dealers: Professional pricing tools"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Enterprise Clients: Bulk processing and API access"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "🏆 Competitive Advantage"
        
        p = tf.add_paragraph()
        p.text = "Superior ML accuracy (98.3% vs industry 85-90%)"
        p.level = 1

    def add_technical_architecture(self):
        """Add technical architecture slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Technical Architecture"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🏗️ Modern Microservices Design"
        
        p = tf.add_paragraph()
        p.text = "Frontend (Next.js) ↔ Nginx Load Balancer ↔ Backend (FastAPI)"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "🛠️ Technology Stack"
        
        p = tf.add_paragraph()
        p.text = "• Frontend: Next.js 14, TypeScript, Tailwind CSS"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Backend: FastAPI, Python, CatBoost ML"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Infrastructure: Docker, Nginx, Ubuntu Server"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Database: MongoDB with Mongoose ODM"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Security: SSL/HTTPS, Rate Limiting, CORS"
        p.level = 1

    def add_ml_excellence(self):
        """Add machine learning excellence slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Machine Learning Excellence"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🧠 CatBoost Model Performance"
        
        p = tf.add_paragraph()
        p.text = "• Accuracy: 98.3% prediction precision"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Features: 87+ vehicle attributes analyzed"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Speed: Sub-second prediction processing"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "📈 Data Analysis Categories"
        
        p = tf.add_paragraph()
        p.text = "• Vehicle Specs: Make, model, year, mileage, engine"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Condition Factors: Accidents, damage, ownership history"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Market Data: Regional trends, seasonal adjustments"
        p.level = 1

    def add_pricing_implementation(self):
        """Add pricing page implementation slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Recent Pricing Page Implementation"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "💰 Subscription Tiers"
        
        p = tf.add_paragraph()
        p.text = "• Basic: $9.99/month - 10 predictions, basic features"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Premium: $24.99/month - Unlimited, advanced analytics"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Enterprise: $99.99/month - Custom solutions, API access"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "💳 Payment Integration"
        
        p = tf.add_paragraph()
        p.text = "• Stripe: Credit cards, digital wallets, international"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• PayPal: Account-based payments, buyer protection"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "⏰ 'Coming Soon' Strategy"
        
        p = tf.add_paragraph()
        p.text = "Professional interface with clear future availability indicators"
        p.level = 1

    def add_security_infrastructure(self):
        """Add security and infrastructure slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Security & Infrastructure"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🔒 Enterprise-Grade Security"
        
        p = tf.add_paragraph()
        p.text = "• SSL/HTTPS: Let's Encrypt with auto-renewal"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Rate Limiting: API protection (10 req/s API, 30 req/s general)"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Authentication: NextAuth.js with OAuth integration"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "🏗️ Production Infrastructure"
        
        p = tf.add_paragraph()
        p.text = "• Server: CarUbuntu (157.245.142.242)"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Containerization: Docker with health monitoring"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Load Balancing: Nginx reverse proxy"
        p.level = 1

    def add_business_model(self):
        """Add business model slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Business Model & Revenue Strategy"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "💼 SaaS Revenue Model"
        
        p = tf.add_paragraph()
        p.text = "• Subscription Tiers: Basic, Premium, Enterprise"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• API Monetization: Usage-based pricing"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Enterprise Solutions: Custom pricing and features"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "📊 Market Opportunity"
        
        p = tf.add_paragraph()
        p.text = "• Automotive Industry: $2.7 trillion global market"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Digital Transformation: Growing demand for AI solutions"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• B2B Potential: Dealer networks and enterprise clients"
        p.level = 1

    def add_future_roadmap(self):
        """Add future roadmap slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Future Development Roadmap"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🚀 Phase 1: Payment Integration (Q3 2025)"
        
        p = tf.add_paragraph()
        p.text = "Complete Stripe/PayPal implementation, subscription management"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "📱 Phase 2: Mobile & API Expansion (Q4 2025)"
        
        p = tf.add_paragraph()
        p.text = "React Native app, public API, advanced analytics"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "🌐 Phase 3: Scale & Enterprise (Q1 2026)"
        
        p = tf.add_paragraph()
        p.text = "Kubernetes deployment, international expansion"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "🔮 Future Innovations"
        
        p = tf.add_paragraph()
        p.text = "Computer vision, blockchain integration, IoT connectivity"
        p.level = 1

    def add_conclusion(self):
        """Add conclusion slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Conclusion & Next Steps"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🎯 Strategic Position"
        
        p = tf.add_paragraph()
        p.text = "Market-leading solution combining technical excellence with business readiness"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "🚀 Immediate Priorities"
        
        p = tf.add_paragraph()
        p.text = "1. Payment Integration: Complete Stripe/PayPal implementation"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "2. User Acquisition: Launch marketing and SEO campaigns"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "3. API Monetization: Activate developer program"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "💡 Investment Opportunity"
        
        p = tf.add_paragraph()
        p.text = "Proven technology with demonstrated accuracy and clear revenue model"
        p.level = 1

    def generate_presentation(self):
        """Generate the complete PowerPoint presentation"""
        print("Generating Car Price Predictor PowerPoint Presentation...")
        
        # Add all slides
        self.add_title_slide()
        self.add_executive_summary()
        self.add_platform_overview()
        self.add_technical_architecture()
        self.add_ml_excellence()
        self.add_pricing_implementation()
        self.add_security_infrastructure()
        self.add_business_model()
        self.add_future_roadmap()
        self.add_conclusion()
        
        # Save the presentation
        self.prs.save(self.filename)
        print(f"PowerPoint presentation generated successfully: {self.filename}")

if __name__ == "__main__":
    generator = PowerPointGenerator()
    generator.generate_presentation()
